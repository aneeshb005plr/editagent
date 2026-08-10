"""
app/documents/pptx_parser.py

Parses .pptx files into the common ParsedDocument representation
(app/documents/base.py). Built and verified against python-pptx 1.0.2
- API specifics confirmed against the real library (title-shape
detection via shape.is_placeholder + placeholder_format.type ==
PP_PLACEHOLDER.TITLE, image bytes via shape.image.blob, chart labels
via chart.chart_title / chart.series / chart.plots[0].categories -
all verified against a generated fixture, not assumed).

SMARTART DETECTION IS UNVERIFIED BY FIXTURE TEST, STILL: python-pptx
has no API to CREATE SmartArt diagrams, so no fixture could be built
to prove the detection path works end-to-end - this remains true
across every review round so far. Implemented per the documented
OOXML namespace for diagrams
(http://schemas.openxmlformats.org/drawingml/2006/diagram) - this
MUST be confirmed against a real SmartArt-containing deck during the
large-file spike before being trusted in production. The confident
tone of _is_smartart's ordering-safety docstring applies ONLY to the
ordering logic (which IS verified) - it does not mean the diagram-
uri match itself has been proven against a real file.

GROUPED SHAPES: recursed into (a group can contain pictures/tables/
text), one level of real-world complexity handled; deeply nested
groups are walked recursively with no depth limit - acceptable
for now, revisit only if a pathological deck causes problems.

NOT YET LOAD-TESTED AT 100MB.
"""

from __future__ import annotations

import logging

from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.media import SPEAKER_IMAGE_BYTES
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture

# python-pptx's shape base class (pptx.shapes.base.BaseShape) is not
# re-exported from pptx.shapes' top level in 1.0.2 - confirmed by
# import error during smoke testing, not assumed. Using `Any` for
# these type hints rather than chasing the internal module path,
# since it's documentation-only here and adds no runtime behavior.
BaseShape = Any

from app.documents.base import (
    ContentBlock,
    ContentKind,
    ExtractionStatus,
    Location,
    ParsedDocument,
    UnsupportedItem,
    UnsupportedKind,
)

logger = logging.getLogger("app.documents.pptx_parser")

_DIAGRAM_NS_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def _is_title_shape(shape: BaseShape) -> bool:
    if not shape.is_placeholder:
        return False
    ph_type = shape.placeholder_format.type
    return ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


def _is_smartart(shape: BaseShape) -> bool:
    """SmartArt is always a graphicFrame whose graphicData uri points
    at the OOXML diagram namespace. Gated on isinstance(shape,
    GraphicFrame) FIRST (not just shape_type != CHART) - the prior
    version ran an unbounded `.//` XML subtree search on every single
    non-chart shape, which is real, avoidable cost at 100MB/many-
    shapes scale. python-pptx doesn't expose a has_smartart flag, so
    the underlying XML still has to be checked directly - see module
    docstring re: unverified-by-fixture status for the diagram-uri
    match itself.

    ORDERING NOTE, confirmed deliberate, not accidental: this runs
    BEFORE the has_chart/has_table checks in _process_shape. A
    populated table/chart placeholder is a GraphicFrame subclass too
    (PlaceholderGraphicFrame), so it also passes the isinstance()
    gate here - but a chart's graphicData uri is the chart namespace
    and a table's is the table namespace, neither of which collides
    with the diagram namespace checked below. So this function
    correctly returns False for populated chart/table placeholders
    and control correctly falls through to the has_chart/has_table
    checks. Do not reorder without re-verifying this non-collision
    still holds."""

    if not isinstance(shape, GraphicFrame):
        return False

    try:
        graphic_data = shape._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData"
        )
        if graphic_data is not None:
            uri = graphic_data.get("uri", "")
            return uri == _DIAGRAM_NS_URI
    except AttributeError:
        return False
    return False


def _extract_chart_labels(shape: BaseShape, location: Location) -> list[ContentBlock]:
    """Chart title/series names/category labels are read directly via
    the chart's own data API - NOT a vision call. The chart's actual
    plotted DATA and visual layout remain unreviewed (flagged
    separately as UnsupportedKind.CHART) - see base.py docstring."""

    blocks: list[ContentBlock] = []
    chart = shape.chart

    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text:
                blocks.append(
                    ContentBlock(
                        text=title_text,
                        kind=ContentKind.CHART_LABEL,
                        location=location,
                    )
                )

        for series in chart.series:
            if series.name and series.name.strip():
                blocks.append(
                    ContentBlock(
                        text=series.name.strip(),
                        kind=ContentKind.CHART_LABEL,
                        location=location,
                    )
                )
        for category in chart.plots[0].categories:
            cat_text = str(category).strip()
            if cat_text:
                blocks.append(
                    ContentBlock(
                        text=cat_text,
                        kind=ContentKind.CHART_LABEL,
                        location=location,
                    )
                )
    except (IndexError, AttributeError):
        # Covers title extraction too now, not just series/categories -
        # a malformed chart_title.text_frame previously raised
        # uncaught since only the series/category block was guarded.
        logger.warning(
            "Could not read title/series/category labels for chart at %s",
            location.display(),
            exc_info=True,
        )

    return blocks


def _process_shape(
    shape: BaseShape,
    slide_number: int,
    parsed: ParsedDocument,
    table_counter: list[int],
) -> None:
    """table_counter is a 1-item mutable list acting as a per-slide
    counter, passed by reference through the recursion (including
    into groups) so multiple tables on one slide get distinct
    table_index values instead of colliding on Row 0, Col 0 - same
    class of bug caught and fixed in docx_parser's table handling."""

    location = Location(slide_number=slide_number)

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            _process_shape(sub_shape, slide_number, parsed, table_counter)
        return

    if _is_smartart(shape):
        parsed.unsupported_items.append(
            UnsupportedItem(
                kind=UnsupportedKind.SMARTART,
                location=location,
                note="SmartArt diagram - structure/content not reviewed in current scope",
                extraction_status=ExtractionStatus.NOT_APPLICABLE,
            )
        )
        return

    if getattr(shape, "has_chart", False):
        parsed.blocks.extend(_extract_chart_labels(shape, location))
        parsed.unsupported_items.append(
            UnsupportedItem(
                kind=UnsupportedKind.CHART,
                location=location,
                note="chart data/visual layout not reviewed - labels extracted separately",
                extraction_status=ExtractionStatus.NOT_APPLICABLE,
            )
        )
        return

    if getattr(shape, "has_table", False):
        this_table_index = table_counter[0]
        table_counter[0] += 1
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.is_spanned:
                    # Explicit skip, not relying on empty-text
                    # filtering as an implicit side effect - confirmed
                    # by testing that a spanned cell's .text is
                    # already empty (python-pptx does not duplicate
                    # the merged-origin cell's text into spanned grid
                    # positions), but making this deliberate rather
                    # than incidental.
                    continue
                cell_text = cell.text.strip()
                if cell_text:
                    parsed.blocks.append(
                        ContentBlock(
                            text=cell_text,
                            kind=ContentKind.TABLE_CELL,
                            location=Location(
                                slide_number=slide_number,
                                table_index=this_table_index,
                                row_index=row_idx,
                                column_index=col_idx,
                            ),
                        )
                    )
        return

    is_picture = (
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE or isinstance(shape, Picture)
    )
    # isinstance() check is the load-bearing one: a POPULATED picture
    # placeholder (python-pptx's PlaceholderPicture) is a genuine
    # Picture subclass with a working .image.blob, but its shape_type
    # reports MSO_SHAPE_TYPE.PLACEHOLDER, not PICTURE - reproduced
    # directly against a real "Picture with Caption" layout during
    # verification. Without the isinstance() check, populated picture
    # placeholders silently vanished entirely (matched none of the
    # branches below, no error, no flag - just gone).
    if is_picture:
        try:
            image = shape.image
            parsed.unsupported_items.append(
                UnsupportedItem(
                    kind=UnsupportedKind.IMAGE,
                    location=location,
                    note=f"content_type={image.content_type}",
                    raw_bytes=image.blob,
                    extraction_status=ExtractionStatus.PENDING,
                )
            )
        except (AttributeError, ValueError):
            # Linked (not embedded) images raise here in python-pptx -
            # flag without bytes rather than crash the parse.
            logger.warning(
                "Could not read image bytes for picture on slide %d "
                "(possibly a linked, not embedded, image)",
                slide_number,
                exc_info=True,
            )
            parsed.unsupported_items.append(
                UnsupportedItem(
                    kind=UnsupportedKind.IMAGE,
                    location=location,
                    note="image bytes unavailable (possibly linked, not embedded)",
                    extraction_status=ExtractionStatus.FAILED,
                )
            )
        return

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            kind = (
                ContentKind.HEADING
                if _is_title_shape(shape)
                else ContentKind.SLIDE_TEXT
            )
            parsed.blocks.append(
                ContentBlock(text=text, kind=kind, location=location)
            )
        return

    if shape.shape_type in (
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
    ):
        parsed.unsupported_items.append(
            UnsupportedItem(
                kind=UnsupportedKind.EMBEDDED_OBJECT,
                location=location,
                note="OLE-embedded/linked object - no extraction path in current scope",
                extraction_status=ExtractionStatus.NOT_APPLICABLE,
            )
        )
        return

    if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
        # Movie and Picture are SIBLINGS under _BasePicture in
        # python-pptx (confirmed via source: class Movie(_BasePicture),
        # class Picture(_BasePicture)) - NOT parent/child - so a movie
        # was never actually caught by the isinstance(shape, Picture)
        # check above, and correctly wasn't silently dropped either
        # (it would have reached the safety net below). But it WAS
        # being mislabeled as generic EMBEDDED_OBJECT, discarding a
        # real, extractable poster-frame image along the way.
        #
        # NOT distinguishing audio vs. video via shape.media_type:
        # confirmed by reading the source that Movie.media_type is
        # HARDCODED to unconditionally return PP_MEDIA_TYPE.MOVIE
        # regardless of the underlying media - it cannot tell audio
        # from video, so gating on it would be a no-op. Instead,
        # comparing the extracted poster_frame bytes directly against
        # pptx.media.SPEAKER_IMAGE_BYTES - confirmed this is literally
        # the exact fallback python-pptx's own add_movie() uses when
        # no real poster frame was supplied (audio clips very commonly
        # hit this path). If it matches, this is a generic stock icon,
        # not real reviewable content - flagged as such rather than
        # sent through vision extraction as if it were meaningful.
        poster_frame = getattr(shape, "poster_frame", None)

        if poster_frame is not None and poster_frame.blob == SPEAKER_IMAGE_BYTES:
            parsed.unsupported_items.append(
                UnsupportedItem(
                    kind=UnsupportedKind.EMBEDDED_OBJECT,
                    location=location,
                    note="media (audio/video) - generic placeholder icon only, no real poster frame to review",
                    extraction_status=ExtractionStatus.NOT_APPLICABLE,
                )
            )
        elif poster_frame is not None:
            parsed.unsupported_items.append(
                UnsupportedItem(
                    kind=UnsupportedKind.IMAGE,
                    location=location,
                    note=f"media poster frame, content_type={poster_frame.content_type}",
                    raw_bytes=poster_frame.blob,
                    extraction_status=ExtractionStatus.PENDING,
                )
            )
        else:
            parsed.unsupported_items.append(
                UnsupportedItem(
                    kind=UnsupportedKind.EMBEDDED_OBJECT,
                    location=location,
                    note="media (audio/video) - no poster frame available",
                    extraction_status=ExtractionStatus.NOT_APPLICABLE,
                )
            )
        return

    # SAFETY NET: anything reaching here (connectors and any other
    # shape type not explicitly handled above) would previously
    # vanish with no trace - no block, no unsupported item, no log.
    # Flagging generically instead means nothing is ever silently
    # dropped, and this branch firing during real-file testing is
    # itself a signal that a new shape category needs its own
    # explicit handling. (Media used to be listed here too, but has
    # its own branch above now - it was never actually silently
    # dropped, just mislabeled when it landed here.)
    parsed.unsupported_items.append(
        UnsupportedItem(
            kind=UnsupportedKind.EMBEDDED_OBJECT,
            location=location,
            note=f"unhandled shape type ({shape.shape_type}) - not reviewed",
            extraction_status=ExtractionStatus.NOT_APPLICABLE,
        )
    )


def parse_pptx(file_bytes: bytes, filename: str) -> ParsedDocument:
    """Parses a .pptx file's readable content into a ParsedDocument.

    File-size validation is the DISPATCHER's responsibility (single
    source of truth against settings.MAX_FILE_SIZE_MB) - this
    function assumes it's already been called with an acceptable
    size and focuses purely on content extraction.
    """

    import io

    prs = Presentation(io.BytesIO(file_bytes))
    parsed = ParsedDocument(source_filename=filename, file_type="pptx")

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1  # 1-indexed for user-facing display
        table_counter = [0]  # per-slide, reset each slide - see _process_shape docstring

        for shape in slide.shapes:
            _process_shape(shape, slide_number, parsed, table_counter)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parsed.blocks.append(
                    ContentBlock(
                        text=notes_text,
                        kind=ContentKind.SLIDE_NOTES,
                        location=Location(slide_number=slide_number),
                    )
                )

    logger.info(
        "Parsed %s: %d text blocks, %d unsupported items (%.0f chars)",
        filename,
        len(parsed.blocks),
        len(parsed.unsupported_items),
        parsed.total_char_count,
    )

    return parsed